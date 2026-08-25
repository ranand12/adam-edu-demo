from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT
BLANK = prs.slide_layouts[6]

BLUE = RGBColor(0x42, 0x85, 0xF4)
RED = RGBColor(0xEA, 0x43, 0x35)
YELLOW = RGBColor(0xFB, 0xBC, 0x04)
GREEN = RGBColor(0x34, 0xA8, 0x53)
DARK_BLUE = RGBColor(0x18, 0x5A, 0xBC)
PURPLE = RGBColor(0x8E, 0x24, 0xAA)
TEAL = RGBColor(0x00, 0x89, 0x7B)
CORAL = RGBColor(0xEE, 0x4D, 0x5D)
SKY = RGBColor(0x07, 0x8E, 0xFB)
TEXT_COLOR = RGBColor(0x20, 0x21, 0x24)
SECONDARY = RGBColor(0x5F, 0x63, 0x68)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SURFACE = RGBColor(0xF8, 0xF9, 0xFA)
BLUE_200 = RGBColor(0x8A, 0xB4, 0xF8)


def add_banner(slide):
    band_w = SLIDE_WIDTH // 4
    h = Emu(Inches(0.1).emu)
    y = SLIDE_HEIGHT - h
    for i, c in enumerate([BLUE, RED, YELLOW, GREEN]):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, band_w * i, y, band_w, h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = c
        shape.line.fill.background()


def add_logo(slide, dark=False):
    txBox = slide.shapes.add_textbox(Inches(0.5), SLIDE_HEIGHT - Inches(0.6), Inches(2), Inches(0.3))
    p = txBox.text_frame.paragraphs[0]
    for ch, color in [('G', BLUE), ('o', RED), ('o', YELLOW), ('g', BLUE), ('l', GREEN), ('e', RED)]:
        run = p.add_run()
        run.text = ch
        run.font.size = Pt(14)
        run.font.color.rgb = color
        run.font.name = 'Arial'
    run = p.add_run()
    run.text = ' Cloud'
    run.font.size = Pt(14)
    run.font.color.rgb = RGBColor(0x99, 0x99, 0x99) if dark else SECONDARY
    run.font.name = 'Arial'


def add_text(slide, left, top, width, height, text, size=18, bold=False, color=TEXT_COLOR, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = 'Arial'
    return tf


def add_rich_text(slide, left, top, width, height, segments, size=18, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    for text, color, bold in segments:
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = 'Arial'
    return tf


def add_card(slide, left, top, width, height, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = SURFACE
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(2.5)
    else:
        shape.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
        shape.line.width = Pt(1)
    return shape


def add_badge(slide, left, top, color, size=Inches(0.6)):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


# ── Slide 1: Title ──
s = prs.slides.add_slide(BLANK)
add_logo(s)
add_rich_text(s, Inches(0.8), Inches(2.0), Inches(10), Inches(3), [
    ('What Can\n', TEXT_COLOR, True),
    ('Gemini Enterprise', BLUE, True),
    ('\nDo?', TEXT_COLOR, True),
], size=44)
add_text(s, Inches(0.8), Inches(5.2), Inches(8), Inches(0.5),
         '9 capabilities transforming higher education', size=18, color=SECONDARY)
add_banner(s)

# ── Slide 2: Content & Research ──
s = prs.slides.add_slide(BLANK)
add_logo(s)
add_text(s, Inches(0.8), Inches(0.5), Inches(4), Inches(0.3),
         'CONTENT & RESEARCH', size=11, bold=True, color=BLUE)

add_rich_text(s, Inches(0.8), Inches(0.85), Inches(10), Inches(0.6), [
    ('Analyze, Search, and ', TEXT_COLOR, True),
    ('Create', BLUE, True),
], size=30)

# Card 1: AI Assistant
card_y = Inches(1.8)
card_h = Inches(1.6)
card_w = Inches(11.5)
add_card(s, Inches(0.8), card_y, card_w, card_h, BLUE)
add_badge(s, Inches(1.1), card_y + Inches(0.3), BLUE)
add_text(s, Inches(1.9), card_y + Inches(0.25), Inches(9), Inches(0.4),
         'AI Assistant — File Analysis', size=20, bold=True)
add_text(s, Inches(1.9), card_y + Inches(0.7), Inches(9.5), Inches(0.8),
         'Upload research papers, datasets, and documents. Gemini extracts key insights, summarizes findings, and answers questions about your files.',
         size=14, color=SECONDARY)

# Card 2: Grounded Web Search
card_y2 = Inches(3.6)
add_card(s, Inches(0.8), card_y2, card_w, card_h, RED)
add_badge(s, Inches(1.1), card_y2 + Inches(0.3), RED)
add_text(s, Inches(1.9), card_y2 + Inches(0.25), Inches(9), Inches(0.4),
         'Grounded Web Search', size=20, bold=True)
add_text(s, Inches(1.9), card_y2 + Inches(0.7), Inches(9.5), Inches(0.8),
         'Search the live web for the latest literature, grant opportunities, and research developments — with cited sources you can verify.',
         size=14, color=SECONDARY)

# Card 3: Media Generation
card_y3 = Inches(5.4)
add_card(s, Inches(0.8), card_y3, card_w, card_h, YELLOW)
add_badge(s, Inches(1.1), card_y3 + Inches(0.3), YELLOW)
add_text(s, Inches(1.9), card_y3 + Inches(0.25), Inches(9), Inches(0.4),
         'Media Generation', size=20, bold=True)
add_text(s, Inches(1.9), card_y3 + Inches(0.7), Inches(9.5), Inches(0.8),
         'Create diagrams, research posters, infographics, and presentation visuals directly from prompts — no design tools required.',
         size=14, color=SECONDARY)
add_banner(s)

# ── Slide 3: Knowledge & Discovery ──
s = prs.slides.add_slide(BLANK)
add_logo(s)
add_text(s, Inches(0.8), Inches(0.5), Inches(4), Inches(0.3),
         'KNOWLEDGE & DISCOVERY', size=11, bold=True, color=GREEN)

add_rich_text(s, Inches(0.8), Inches(0.85), Inches(10), Inches(0.6), [
    ('Synthesize, Discover, and ', TEXT_COLOR, True),
    ('Explore', BLUE, True),
], size=30)

# Card 1: Research Knowledge Agents
add_card(s, Inches(0.8), card_y, card_w, card_h, GREEN)
add_badge(s, Inches(1.1), card_y + Inches(0.3), GREEN)
add_text(s, Inches(1.9), card_y + Inches(0.25), Inches(9), Inches(0.4),
         'Research Knowledge Agents', size=20, bold=True)
add_text(s, Inches(1.9), card_y + Inches(0.7), Inches(9.5), Inches(0.8),
         'Build reusable Gems — custom AI assistants for lab protocols, grant writing, methodology guidance, and domain-specific expertise.',
         size=14, color=SECONDARY)

# Card 2: NotebookLM
add_card(s, Inches(0.8), card_y2, card_w, card_h, PURPLE)
add_badge(s, Inches(1.1), card_y2 + Inches(0.3), PURPLE)
add_text(s, Inches(1.9), card_y2 + Inches(0.25), Inches(9), Inches(0.4),
         'NotebookLM Literature Review', size=20, bold=True)
add_text(s, Inches(1.9), card_y2 + Inches(0.7), Inches(9.5), Inches(0.8),
         'Upload multiple papers and let NotebookLM synthesize themes, generate audio overviews, and answer cross-document questions.',
         size=14, color=SECONDARY)

# Card 3: Deep Research
add_card(s, Inches(0.8), card_y3, card_w, card_h, TEAL)
add_badge(s, Inches(1.1), card_y3 + Inches(0.3), TEAL)
add_text(s, Inches(1.9), card_y3 + Inches(0.25), Inches(9), Inches(0.4),
         'Deep Research', size=20, bold=True)
add_text(s, Inches(1.9), card_y3 + Inches(0.7), Inches(9.5), Inches(0.8),
         'Launch comprehensive, web-wide research queries that return structured reports with citations — the depth of a research assistant in minutes.',
         size=14, color=SECONDARY)
add_banner(s)

# ── Slide 4: Platform & Extensibility ──
s = prs.slides.add_slide(BLANK)
add_logo(s)
add_text(s, Inches(0.8), Inches(0.5), Inches(4), Inches(0.3),
         'PLATFORM & EXTENSIBILITY', size=11, bold=True, color=DARK_BLUE)

add_rich_text(s, Inches(0.8), Inches(0.85), Inches(10), Inches(0.6), [
    ('Connect, Extend, and ', TEXT_COLOR, True),
    ('Build', BLUE, True),
], size=30)

# Card 1: Multi-Model Access
add_card(s, Inches(0.8), card_y, card_w, card_h, DARK_BLUE)
add_badge(s, Inches(1.1), card_y + Inches(0.3), DARK_BLUE)
add_text(s, Inches(1.9), card_y + Inches(0.25), Inches(9), Inches(0.4),
         'Multi-Model Access', size=20, bold=True)
add_text(s, Inches(1.9), card_y + Inches(0.7), Inches(9.5), Inches(0.8),
         'Connect to multiple AI models — including Claude from Anthropic — directly through Gemini Enterprise. Choose the best model for every task.',
         size=14, color=SECONDARY)

# Card 2: External Connectors
add_card(s, Inches(0.8), card_y2, card_w, card_h, SKY)
add_badge(s, Inches(1.1), card_y2 + Inches(0.3), SKY)
add_text(s, Inches(1.9), card_y2 + Inches(0.25), Inches(9), Inches(0.4),
         'External Connectors', size=20, bold=True)
add_text(s, Inches(1.9), card_y2 + Inches(0.7), Inches(9.5), Inches(0.8),
         "Link Gemini to external data sources, APIs, and institutional tools — bring your university's data into the AI conversation.",
         size=14, color=SECONDARY)

# Card 3: Custom Agents
add_card(s, Inches(0.8), card_y3, card_w, card_h, CORAL)
add_badge(s, Inches(1.1), card_y3 + Inches(0.3), CORAL)
add_text(s, Inches(1.9), card_y3 + Inches(0.25), Inches(9), Inches(0.4),
         'Custom Agents', size=20, bold=True)
add_rich_text(s, Inches(1.9), card_y3 + Inches(0.7), Inches(9.5), Inches(0.8), [
    ('Build purpose-built AI agents for specific university needs — like ', SECONDARY, False),
    ('Grants AI', BLUE, True),
    (' for funding discovery and ', SECONDARY, False),
    ('AI Tutor', BLUE, True),
    (' for personalized student support.', SECONDARY, False),
], size=14)
add_banner(s)

# ── Slide 5: Demos ──
s = prs.slides.add_slide(BLANK)
add_logo(s)
add_rich_text(s, Inches(0.8), Inches(2.8), Inches(11), Inches(1.5), [
    ('Demos', BLUE, True),
], size=56, align=PP_ALIGN.LEFT)
add_banner(s)

# ── Slide 6: Thank You ──
s = prs.slides.add_slide(BLANK)
add_logo(s)
add_rich_text(s, Inches(0.8), Inches(2.8), Inches(11), Inches(1.5), [
    ('Thank ', TEXT_COLOR, True),
    ('you', BLUE, True),
], size=56, align=PP_ALIGN.LEFT)
add_banner(s)

out = '/home/user/AI Projects/Adam EDU Demo/adam-edu-demo-slides.pptx'
prs.save(out)
print(f'Saved: {out}')
